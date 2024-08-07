from pptx import Presentation

from chrisbase.data import *
from chrisbase.io import *
from chrisbase.util import *

logger = logging.getLogger(__name__)
args = CommonArguments(
    env=ProjectEnv(
        project="SlideScanner",
        job_name=Path(__file__).stem,
        msg_level=logging.INFO,
        msg_format=LoggingFormat.PRINT_00,
    )
)
args.info_args()

BR_MARK = "<BR>"


def remove_all_slides(prs: Presentation):
    for i, x in enumerate(prs.slides._sldIdLst):
        prs.part.drop_rel(x.rId)
        del prs.slides._sldIdLst[i]
    return prs


def get_shape_text(shape):
    return ' '.join(BR_MARK.join(x.strip() for x in str(shape.text).strip().replace("\t", " ").replace("\r", "").replace("\x0b", "\n").split("\n")).split())


def check_shape_name(file_path):
    file_cont = list()
    for slide in Presentation(file_path).slides:
        slide_cont = []
        for shape in slide.shapes:
            slide_cont.append(shape.name)
        if (slide_cont == ["제목 2", "내용 개체 틀 3"] or
                slide_cont == ["제목 1", "내용 개체 틀 2"]):
            pass
        else:
            file_cont.append(slide_cont)
    return file_cont


index_pattern = re.compile(r"\([^\)]+?[0-9]+?\)")
acceptable_index_prefix = ['기도회', '찬']

def scan_pptx(file_path):
    title_text = None
    page_texts = list()

    def index_prefix(x):
        return [a for a in re.split(r"[0-9]+|\(|\)", x) if a][0]

    def to_fname(x):
        return ''.join(re.compile('[ㄱ-ㅎ가-힣A-Za-z0-9]+').findall(x))

    index = [x for x in index_pattern.findall(Path(file_path).stem)
             if index_prefix(x) in acceptable_index_prefix]

    for slide in Presentation(file_path).slides:
        assert len(slide.shapes) == 2
        first_shape = slide.shapes[0]
        assert first_shape.name in [
            "제목 1",
            "제목 2",
            "Title 1",
        ], f"first_shape.name={first_shape.name}"
        second_shape = slide.shapes[1]
        assert second_shape.name in [
            "내용 개체 틀 2",
            "내용 개체 틀 3",
            "Text Placeholder 2",
            "Content Placeholder 2",
        ], f"second_shape.name={second_shape.name}"
        if not title_text:
            title_text = get_shape_text(first_shape)
        page_texts.append(get_shape_text(second_shape))
    if title_text and page_texts:
        lines_ns = to_fname(''.join(''.join(page_texts[0].split(BR_MARK)[:2]).split()))
        title_ns = to_fname(''.join(title_text.split("(")[0].split()))
        index.append(lines_ns if lines_ns.startswith(title_ns) else f"{lines_ns} ({title_ns})")
        return {
            "fname": ' '.join(index),
            "title": title_text,
            "pages": page_texts,
        }
    return None


base_path = "resource/base-key.pptx"
input_dir = "/Users/chris/Seafile/love/찬양 PPT"
output_dir = "/Users/chris/Seafile/love/찬양 PPT2"
output_dir = make_dir(output_dir)

with JobTimer(args.env.job_name, rt=1, rb=1, rw=80, rc='=', verbose=1):
    contents_set = set()
    for file in sorted(Path(input_dir).glob("*.pptx")):
        contents = scan_pptx(file)
        if contents:
            contents_set.add(json.dumps(contents, ensure_ascii=False))
    contents_set = [json.loads(x) for x in contents_set]
    contents_set = sorted(contents_set, key=lambda x: x["fname"])
    for contents in contents_set:
        print(json.dumps(contents, indent=4, ensure_ascii=False))
        output_path = output_dir / (contents["fname"] + ".pptx")
        base_pptx = remove_all_slides(Presentation(base_path))
        slide_layout = base_pptx.slide_layouts[0]
        title = contents["title"]
        for page in contents["pages"]:
            slide = base_pptx.slides.add_slide(slide_layout)
            slide.shapes.placeholders[0].text = title
            slide.shapes.placeholders[1].text = page.replace("<BR>", "\n")
        base_pptx.save(output_path)
